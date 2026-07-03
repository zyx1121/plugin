# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx==1.0.2"]
# ///
"""Inspect a WinLab template + an example deck to extract reusable conventions."""
import sys
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def cm(v):
    return round(Emu(v).cm, 2) if v is not None else None


def pt(v):
    return round(v / 12700, 1) if v is not None else None  # EMU->pt for font size (Pt stored in centi-pt actually)


def theme_info(path):
    z = zipfile.ZipFile(path)
    name = next((n for n in z.namelist() if n.startswith("ppt/theme/theme1")), None)
    if not name:
        return "(no theme1)"
    root = etree.fromstring(z.read(name))
    out = ["  clrScheme:"]
    clr = root.find(f".//{{{A}}}clrScheme")
    for child in clr:
        tag = etree.QName(child).localname
        srgb = child.find(f"{{{A}}}srgbClr")
        sysc = child.find(f"{{{A}}}sysClr")
        val = srgb.get("val") if srgb is not None else (sysc.get("lastClr") if sysc is not None else "?")
        out.append(f"    {tag:10} #{val}")
    fs = root.find(f".//{{{A}}}fontScheme")
    out.append("  fontScheme:")
    for major in ("majorFont", "minorFont"):
        mf = fs.find(f"{{{A}}}{major}")
        latin = mf.find(f"{{{A}}}latin")
        ea = mf.find(f"{{{A}}}ea")
        out.append(f"    {major:10} latin={latin.get('typeface')!r:20} ea={ea.get('typeface')!r}")
    return "\n".join(out)


def shape_kind(sp):
    tag = etree.QName(sp.element).localname
    return {"sp": "shape", "cxnSp": "CONNECTOR", "pic": "PIC", "grpSp": "GROUP", "graphicFrame": "gframe"}.get(tag, tag)


def first_text(sp, n=70):
    if not sp.has_text_frame:
        return ""
    t = sp.text_frame.text.replace("\n", " ⏎ ").strip()
    return (t[:n] + "…") if len(t) > n else t


def font_summary(sp):
    if not sp.has_text_frame:
        return ""
    sizes, colors, fonts = set(), set(), set()
    for para in sp.text_frame.paragraphs:
        for run in para.runs:
            f = run.font
            if f.size:
                sizes.add(round(f.size.pt, 1))
            try:
                if f.color and f.color.type is not None and f.color.rgb:
                    colors.add(str(f.color.rgb))
            except Exception:
                pass
            if f.name:
                fonts.add(f.name)
    bits = []
    if sizes:
        bits.append(f"sz={sorted(sizes)}")
    if colors:
        bits.append(f"clr={sorted(colors)}")
    if fonts:
        bits.append(f"font={sorted(fonts)}")
    return " ".join(bits)


def autoshape_detail(sp):
    el = sp.element
    geom = el.find(f".//{{{A}}}prstGeom")
    prst = geom.get("prst") if geom is not None else "?"
    fill = el.find(f".//{{{A}}}spPr/{{{A}}}solidFill/{{{A}}}srgbClr")
    fillv = f"#{fill.get('val')}" if fill is not None else "(theme/none)"
    ln = el.find(f".//{{{A}}}spPr/{{{A}}}ln")
    lnv = ""
    if ln is not None:
        lc = ln.find(f".//{{{A}}}srgbClr")
        lnv = f"line=#{lc.get('val') if lc is not None else '?'} w={ln.get('w')}"
    return f"prst={prst} fill={fillv} {lnv}"


def connector_detail(sp):
    el = sp.element
    st = el.find(f".//{{{A}}}stCxn")
    en = el.find(f".//{{{A}}}endCxn")
    s = f"st(id={st.get('id')},idx={st.get('idx')})" if st is not None else "st(free)"
    e = f"end(id={en.get('id')},idx={en.get('idx')})" if en is not None else "end(free)"
    tail = el.find(f".//{{{A}}}tailEnd")
    arr = f" arrow={tail.get('type')}" if tail is not None else ""
    geom = el.find(f".//{{{A}}}prstGeom")
    return f"{geom.get('prst') if geom is not None else '?'} {s}->{e}{arr}"


def dump_template(path):
    print("=" * 70)
    print(f"TEMPLATE: {path}")
    print("=" * 70)
    prs = Presentation(path)
    print(f"slide size: {prs.slide_width}x{prs.slide_height} EMU = {cm(prs.slide_width)}x{cm(prs.slide_height)} cm  ratio={round(prs.slide_width/prs.slide_height,3)}")
    print(theme_info(path))
    print(f"\nLAYOUTS ({len(prs.slide_layouts)}):")
    for i, lay in enumerate(prs.slide_layouts):
        phs = []
        for ph in lay.placeholders:
            pf = ph.placeholder_format
            phs.append(f"idx={pf.idx} {str(pf.type)} '{ph.name}' @({cm(ph.left)},{cm(ph.top)}) {cm(ph.width)}x{cm(ph.height)}cm")
        print(f"  [{i}] name={lay.name!r}  ph={len(phs)}")
        for p in phs:
            print(f"        {p}")
    print(f"\nTEMPLATE pre-built slides ({len(prs.slides)}):")
    for i, sl in enumerate(prs.slides):
        print(f"  slide{i+1}: layout={sl.slide_layout.name!r} shapes={len(sl.shapes)}")
        for sp in sl.shapes:
            txt = first_text(sp)
            print(f"      {shape_kind(sp):9} '{sp.name}' {('ph'+str(sp.placeholder_format.idx)) if sp.is_placeholder else ''} {('| '+txt) if txt else ''} {font_summary(sp)}")


def dump_deck(path):
    print("\n" + "=" * 70)
    print(f"EXAMPLE DECK: {path}")
    print("=" * 70)
    prs = Presentation(path)
    print(f"slide size: {cm(prs.slide_width)}x{cm(prs.slide_height)} cm  pages={len(prs.slides)}")
    diagram_pages = []
    for i, sl in enumerate(prs.slides):
        kinds = {}
        autos = conns = pics = 0
        for sp in sl.shapes:
            k = shape_kind(sp)
            kinds[k] = kinds.get(k, 0) + 1
            if k == "CONNECTOR":
                conns += 1
            if k == "PIC":
                pics += 1
            el = sp.element
            if k == "shape" and el.find(f".//{{{A}}}prstGeom") is not None and not sp.is_placeholder:
                autos += 1
        title = ""
        for sp in sl.shapes:
            if sp.is_placeholder and sp.placeholder_format.idx == 0:
                title = first_text(sp, 50)
                break
        if not title:
            for sp in sl.shapes:
                t = first_text(sp, 50)
                if t:
                    title = t
                    break
        flag = "  <<< DIAGRAM?" if (autos >= 2 or conns >= 1) else ""
        print(f"  p{i+1:2} layout={sl.slide_layout.name!r:28} {dict(kinds)} auto={autos} conn={conns} | {title}{flag}")
        if autos >= 2 or conns >= 1:
            diagram_pages.append((i + 1, sl))

    print(f"\n--- DIAGRAM PAGE DETAIL ({len(diagram_pages)} pages) ---")
    for pno, sl in diagram_pages:
        print(f"\n  === p{pno} (layout={sl.slide_layout.name!r}) ===")
        for sp in sl.shapes:
            k = shape_kind(sp)
            pos = f"@({cm(sp.left)},{cm(sp.top)}) {cm(sp.width)}x{cm(sp.height)}cm" if sp.left is not None else ""
            if k == "CONNECTOR":
                print(f"      CONNECTOR {pos} {connector_detail(sp)}")
            elif k == "shape" and sp.element.find(f".//{{{A}}}prstGeom") is not None:
                txt = first_text(sp, 40)
                sid = sp.element.find(f".//{{{P}}}cNvPr")
                sidv = sid.get("id") if sid is not None else "?"
                print(f"      AUTOSHAPE id={sidv} {pos} {autoshape_detail(sp)} {('| '+txt) if txt else ''} {font_summary(sp)}")
            elif k == "PIC":
                print(f"      PIC '{sp.name}' {pos}")
            else:
                txt = first_text(sp, 40)
                if txt or k != "shape":
                    print(f"      {k} '{sp.name}' {pos} {('| '+txt) if txt else ''} {font_summary(sp)}")


if __name__ == "__main__":
    dump_template(sys.argv[1])
    dump_deck(sys.argv[2])
