# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx==1.0.2"]
# ///
"""WinLab pptx builder — render a deck JSON spec onto template.pptx.

Usage: uv run builder.py <template.pptx> <spec.json> <out.pptx>

Layouts: cover | outline | section | content | photo | two-col | diagram

Diagram notation is SEMANTIC, not hand-picked (Moody's Physics of Notations:
a visual difference must encode a real semantic difference, else it's noise):
  box  `kind`  -> shape   (shape = "kind of thing")
  edge `kind`  -> line+arrow (solid/strong vs dashed/weak; orthogonal to shape)
A legend is auto-drawn from the kinds actually used (C4: every notation needs a key).
Rounded-vs-square is NOT a semantic channel — all process boxes are one shape.
"""
import json
import sys

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Emu, Pt
from pptx.enum.dml import MSO_THEME_COLOR

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0, 0, 0)
GREY = RGBColor(0x80, 0x80, 0x80)
ACCENT = MSO_THEME_COLOR.ACCENT_1

# citation footnote band (cm) — sits below the content box, above the NYCU footer
CITE_Y, CITE_LEFT, CITE_W = 17.15, 1.5, 30.8

# Master locks latin font to Calibri (no CJK glyphs). For Chinese teaching decks
# we add an East Asian typeface per run so CJK renders consistently on Office;
# latin/numbers still inherit Calibri (brand untouched). Broadly available on Win+Mac Office.
EA_FONT = "Microsoft JhengHei"

# template placeholder idx contract (from inspect)
PH_TITLE, PH_BODY, PH_NUM, PH_SECOND = 0, 1, 2, 21
LAYOUT = {
    "cover": "Title",
    "outline": "Title & Bullets",
    "section": "Section",
    "content": "Title & Bullets",
    "photo": "Title, Bullets & Photo",
    "two-col": "Two Columns",
    "diagram": "Title & Bullets",
}

# box kind -> shape. shape encodes KIND OF THING (Moody: "differences in shape
# are read as differences in kind"; ISO 5807 flowchart conventions).
KIND_SHAPES = {
    "component": MSO_SHAPE.RECTANGLE,            # process / module (default)
    "store": MSO_SHAPE.CAN,                      # datastore / DB (cylinder)
    "decision": MSO_SHAPE.DIAMOND,               # branch / decision
    "io": MSO_SHAPE.PARALLELOGRAM,               # input / output
    "external": MSO_SHAPE.ROUNDED_RECTANGLE,     # external actor / system boundary
}
KIND_LABELS = {"component": "Component", "store": "Store", "decision": "Decision",
               "io": "I/O", "external": "External"}

# edge kind -> (dash style, arrowhead). line texture & arrowhead are two ORTHOGONAL
# channels (UML). Each kind must be visually distinct (Semiotic Clarity: one
# symbol, one meaning) — solid / dashed / dotted are three discriminable textures.
EDGE_STYLES = {
    "flow": (None, "triangle"),     # solid + filled arrow: runtime data flow / sync call
    "async": ("dash", "arrow"),     # dashed + open arrow: async / event stream
    "dep": ("sysDot", "arrow"),     # dotted + open arrow: depends-on / config (≠ async)
}
EDGE_LABELS = {"flow": "data flow", "async": "async / event", "dep": "depends on"}

# diagram canvas (cm) — below the title band; bottom strip reserved for legend
DG_LEFT, DG_TOP, DG_W, DG_H = 1.5, 3.2, 30.8, 13.6
BOX_W, BOX_H = 3.53, 1.75  # default box size (cm)
LEGEND_Y = DG_TOP + DG_H + 0.25


def layout_by_name(prs, name):
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    raise SystemExit(f"layout {name!r} not in template; have {[l.name for l in prs.slide_layouts]}")


def ph(slide, idx):
    for p in slide.placeholders:
        if p.placeholder_format.idx == idx:
            return p
    return None


def set_bullets(tf, bullets):
    """bullets: list of {text, level, bold} or plain str (level 0)."""
    tf.clear()
    for i, b in enumerate(bullets):
        if isinstance(b, str):
            text, level, bold = b, 0, False
        else:
            text, level, bold = b.get("text", ""), b.get("level", 0), b.get("bold", False)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        if bold:
            for r in p.runs:
                r.font.bold = True


# ---------- slide renderers ----------
def render_cover(prs, s):
    slide = prs.slides.add_slide(layout_by_name(prs, LAYOUT["cover"]))
    ph(slide, PH_TITLE).text = s["title"]
    set_bullets(ph(slide, PH_BODY).text_frame, [s.get("date", ""), s.get("author", "詹詠翔")])


def render_outline(prs, s):
    slide = prs.slides.add_slide(layout_by_name(prs, LAYOUT["outline"]))
    ph(slide, PH_TITLE).text = s.get("title", "Outline")
    cur = s.get("current")
    items = [{"text": it, "level": 0, "bold": (i == cur)} for i, it in enumerate(s["items"])]
    set_bullets(ph(slide, PH_BODY).text_frame, items)


def render_section(prs, s):
    slide = prs.slides.add_slide(layout_by_name(prs, LAYOUT["section"]))
    ph(slide, PH_TITLE).text = s["title"]


def render_content(prs, s):
    slide = prs.slides.add_slide(layout_by_name(prs, LAYOUT["content"]))
    ph(slide, PH_TITLE).text = s["title"]
    set_bullets(ph(slide, PH_BODY).text_frame, s.get("bullets", []))


def render_two_col(prs, s):
    slide = prs.slides.add_slide(layout_by_name(prs, LAYOUT["two-col"]))
    ph(slide, PH_TITLE).text = s["title"]
    set_bullets(ph(slide, PH_BODY).text_frame, s.get("left", []))
    second = ph(slide, PH_SECOND)
    if second is not None:
        set_bullets(second.text_frame, s.get("right", []))


# ---------- diagram engine ----------
def _style_box(shape, font_pt):
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.theme_color = ACCENT
    shape.line.width = Pt(2)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(font_pt)
            r.font.color.rgb = BLACK  # autoshape default text is white-on-white otherwise


def _center(geom):
    l, t, w, h = geom
    return (l + w / 2, t + h / 2)


def _sides(a, b):
    """pick connection-point idx (0 top,1 left,2 bottom,3 right) for a->b."""
    acx, acy = _center(a)
    bcx, bcy = _center(b)
    dx, dy = bcx - acx, bcy - acy
    if abs(dx) >= abs(dy):
        return (3, 1) if dx >= 0 else (1, 3)
    return (2, 0) if dy >= 0 else (0, 2)


def _style_edge(conn, kind):
    """line texture + arrowhead from edge kind (solid/strong vs dashed/weak)."""
    dash, arrow = EDGE_STYLES.get(kind, EDGE_STYLES["flow"])
    conn.line.color.theme_color = ACCENT
    conn.line.width = Pt(2)
    ln = conn.line._get_or_add_ln()  # a:ln child order: fill, prstDash, ..., tailEnd
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": arrow, "w": "med", "len": "med"}))


def _aligned(ga, gb):
    ca, cb = _center(ga), _center(gb)
    return abs(ca[0] - cb[0]) < 0.15 or abs(ca[1] - cb[1]) < 0.15


def _draw_legend(slide, box_kinds, edge_kinds, note):
    """auto notation key from the kinds actually used (C4: every diagram needs one)."""
    x = DG_LEFT + 0.2
    for k in box_kinds:
        ms = slide.shapes.add_shape(KIND_SHAPES[k], Cm(x), Cm(LEGEND_Y), Cm(0.7), Cm(0.5))
        ms.fill.solid()
        ms.fill.fore_color.rgb = WHITE
        ms.line.color.theme_color = ACCENT
        ms.line.width = Pt(1.5)
        ms.shadow.inherit = False
        _label(slide, x + 0.85, KIND_LABELS[k])
        x += 0.85 + 2.5
    for k in edge_kinds:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x), Cm(LEGEND_Y + 0.25), Cm(x + 1.0), Cm(LEGEND_Y + 0.25))
        _style_edge(conn, k)
        _label(slide, x + 1.15, EDGE_LABELS[k])
        x += 1.15 + 2.6
    if note:
        _label(slide, x + 0.2, note, width=DG_LEFT + DG_W - x - 0.4)


def _label(slide, x, text, width=2.5):
    lab = slide.shapes.add_textbox(Cm(x), Cm(LEGEND_Y - 0.08), Cm(width), Cm(0.7))
    lab.text_frame.word_wrap = False
    lab.text_frame.text = text
    for r in lab.text_frame.paragraphs[0].runs:
        r.font.size = Pt(10)
        r.font.color.rgb = BLACK


def render_diagram(prs, s):
    slide = prs.slides.add_slide(layout_by_name(prs, LAYOUT["diagram"]))
    ph(slide, PH_TITLE).text = s["title"]
    cols, rows = s.get("cols", 6), s.get("rows", 6)
    cell_w, cell_h = DG_W / cols, DG_H / rows

    # 1. geometry per box id (centered in its grid cell)
    geom = {}
    for b in s["boxes"]:
        bw, bh = b.get("w_cm", BOX_W), b.get("h_cm", BOX_H)
        cx = DG_LEFT + (b["col"] + 0.5) * cell_w
        cy = DG_TOP + (b["row"] + 0.5) * cell_h
        geom[b["id"]] = (cx - bw / 2, cy - bh / 2, bw, bh)

    # 2. zone containers FIRST (z-order behind boxes); dashed = logical boundary
    for z in s.get("zones", []):
        ls = [geom[i] for i in z["boxes"]]
        pad = z.get("pad_cm", 0.45)
        minl, mint = min(g[0] for g in ls) - pad, min(g[1] for g in ls) - pad
        maxr, maxb = max(g[0] + g[2] for g in ls) + pad, max(g[1] + g[3] for g in ls) + pad
        zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(minl), Cm(mint), Cm(maxr - minl), Cm(maxb - mint))
        zone.fill.background()
        zone.line.color.theme_color = ACCENT
        zone.line.width = Pt(2)
        zone.shadow.inherit = False
        zln = zone.line._get_or_add_ln()
        zln.append(zln.makeelement(qn("a:prstDash"), {"val": "dash"}))
        if z.get("label"):  # top-right, white bg masks the border behind it
            lab = slide.shapes.add_textbox(Cm(maxr - 3.5), Cm(mint + 0.12), Cm(3.3), Cm(0.7))
            lab.fill.solid()
            lab.fill.fore_color.rgb = WHITE
            lab.line.fill.background()
            lab.text_frame.word_wrap = False
            lab.text_frame.text = z["label"]
            lab.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            for r in lab.text_frame.paragraphs[0].runs:
                r.font.size = Pt(12)

    # 3. boxes — shape from kind
    shapes = {}
    for b in s["boxes"]:
        l, t, w, h = geom[b["id"]]
        kind = b.get("kind", "component")
        shp = slide.shapes.add_shape(KIND_SHAPES.get(kind, MSO_SHAPE.RECTANGLE), Cm(l), Cm(t), Cm(w), Cm(h))
        shp.text_frame.text = b["text"].replace("\\n", "\n")
        _style_box(shp, b.get("font", 14))
        shapes[b["id"]] = shp

    # 4. glued connectors — line/arrow from kind; straight when aligned
    for e in s.get("edges", []):
        a, b = shapes[e["from"]], shapes[e["to"]]
        sidx, eidx = _sides(geom[e["from"]], geom[e["to"]])
        straight = e.get("straight") or _aligned(geom[e["from"]], geom[e["to"]])
        routing = MSO_CONNECTOR.STRAIGHT if straight else MSO_CONNECTOR.ELBOW
        conn = slide.shapes.add_connector(routing, a.left, a.top, b.left, b.top)
        conn.begin_connect(a, sidx)
        conn.end_connect(b, eidx)
        _style_edge(conn, e.get("kind", "flow"))

    # 5. auto legend from kinds used (preserve first-seen order) + optional note
    box_kinds = list(dict.fromkeys(b.get("kind", "component") for b in s["boxes"]))
    edge_kinds = list(dict.fromkeys(e.get("kind", "flow") for e in s.get("edges", [])))
    _draw_legend(slide, box_kinds, edge_kinds, s.get("note"))


RENDERERS = {
    "cover": render_cover, "outline": render_outline, "section": render_section,
    "content": render_content, "two-col": render_two_col, "diagram": render_diagram,
}


def _apply_ea_font(slide, name):
    """Add an East Asian typeface to every run so CJK renders (master is Calibri-only).
    Run rPr is normally empty/attribute-only, so appending a:ea yields a valid rPr."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rPr = run._r.get_or_add_rPr()
                if rPr.find(qn("a:ea")) is None:
                    rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": name}))


def _set_notes(slide, text):
    """Write speaker notes (Keynote/PPT 的 note = presenter notes, not the diagram footnote)."""
    tf = slide.notes_slide.notes_text_frame
    tf.text = text.replace("\\n", "\n")
    for para in tf.paragraphs:
        for run in para.runs:
            rPr = run._r.get_or_add_rPr()
            if rPr.find(qn("a:ea")) is None:
                rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": EA_FONT}))


def _render_cite(slide, cites):
    """Small grey source line under the content box: '來源：labelA ｜ labelB'.
    cites: list of {label, url}. Labels shown on slide; full URLs go to notes."""
    labels = "　｜　".join(c["label"] for c in cites)
    box = slide.shapes.add_textbox(Cm(CITE_LEFT), Cm(CITE_Y), Cm(CITE_W), Cm(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = "來源：" + labels
    p = tf.paragraphs[0]
    for r in p.runs:
        r.font.size = Pt(11)
        r.font.color.rgb = GREY
        rPr = r._r.get_or_add_rPr()
        if rPr.find(qn("a:ea")) is None:
            rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": EA_FONT}))


def build(template, spec, out):
    prs = Presentation(template)
    # drop pre-built template slides: remove BOTH the sldId ref and the relationship,
    # else orphan slide parts get re-serialized and collide with new ones.
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        prs.part.drop_rel(sld.rId)
        xml_slides.remove(sld)
    for s in spec["slides"]:
        RENDERERS[s["layout"]](prs, s)
        slide = prs.slides[-1]
        _apply_ea_font(slide, EA_FONT)
        cites = s.get("cite") or []
        if cites:
            _render_cite(slide, cites)
        # speaker notes: base text + auto-appended full citation links (slide shows labels only)
        notes = s.get("notes", "")
        if cites:
            links = "\n".join(f"- {c['label']}: {c['url']}" for c in cites)
            notes = (notes + "\n\n參考連結：\n" + links).strip()
        if notes:
            _set_notes(slide, notes)
    prs.save(out)
    print(f"wrote {out} ({len(spec['slides'])} slides)")


if __name__ == "__main__":
    template, spec_path, out = sys.argv[1], sys.argv[2], sys.argv[3]
    build(template, json.load(open(spec_path)), out)
